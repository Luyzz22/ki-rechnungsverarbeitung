from __future__ import annotations

from io import BytesIO
from typing import Any, Optional

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from .data import MBRData
from .types import MBRNarrative


def format_eur(x: float) -> str:
    s = f"{x:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    return f"{s} €"


def _remove_shape(shape) -> None:
    el = shape._element
    el.getparent().remove(el)


def _find_shape_with_token(slide, token: str):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and token in (shape.text or ""):
            return shape
    return None


def _replace_tokens_in_text(slide, mapping: dict[str, str]) -> None:
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        tf = shape.text_frame
        for p in tf.paragraphs:
            for r in p.runs:
                for k, v in mapping.items():
                    if k in r.text:
                        r.text = r.text.replace(k, v)


def _set_bullets(shape, bullets: list[str]) -> None:
    tf = shape.text_frame
    tf.clear()
    if not bullets:
        tf.text = ""
        return
    tf.text = bullets[0]
    for b in bullets[1:]:
        p = tf.add_paragraph()
        p.text = b
        p.level = 0


def _replace_placeholder_with_bullets(slide, placeholder_text: str, bullets: list[str]) -> None:
    """Find a shape containing placeholder_text and replace with bullet list."""
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if placeholder_text in (shape.text or ""):
            tf = shape.text_frame
            tf.clear()
            if bullets:
                tf.text = "• " + bullets[0]
                for b in bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = "• " + b
                    p.level = 0
            else:
                tf.text = "Keine Daten verfügbar."
            return


def _add_table_top_suppliers(slide, placeholder_shape, suppliers: list[tuple[str, float]]) -> None:
    left, top, width, height = placeholder_shape.left, placeholder_shape.top, placeholder_shape.width, placeholder_shape.height
    _remove_shape(placeholder_shape)

    rows = max(2, len(suppliers) + 1)
    cols = 2
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    table.cell(0, 0).text = "Lieferant"
    table.cell(0, 1).text = "Netto"

    for i, (name, amt) in enumerate(suppliers, start=1):
        table.cell(i, 0).text = name
        table.cell(i, 1).text = format_eur(amt)


def _add_budget_chart(slide, placeholder_shape, categories: list[tuple[str, float, float]]) -> None:
    left, top, width, height = placeholder_shape.left, placeholder_shape.top, placeholder_shape.width, placeholder_shape.height
    _remove_shape(placeholder_shape)

    chart_data = CategoryChartData()
    chart_data.categories = [c[0] for c in categories]
    chart_data.add_series("Ist (Netto)", [c[1] for c in categories])
    chart_data.add_series("Budget", [c[2] for c in categories])

    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        left, top, width, height,
        chart_data,
    )


def render_presentation_from_template(
    template_path: str,
    data: MBRData,
    narrative: MBRNarrative,
) -> bytes:
    prs = Presentation(template_path)

    # Basic token map
    token_map = {
        "{{MBR_MONTH}}": narrative.month_label,
        "{{COVERAGE_NOTE}}": data.coverage_note,
        "{{INVOICE_COUNT}}": str(data.invoice_count),
        "{{TOTAL_NET}}": format_eur(data.total_net),
        "{{TOTAL_GROSS}}": format_eur(data.total_gross),
        "{{CLOSING_STATEMENT}}": narrative.closing_statement,
    }

    # Slide-wide token replacement
    for slide in prs.slides:
        _replace_tokens_in_text(slide, token_map)

    # Bullet placeholders
    bullet_slots = [
        ("{{EXEC_SUMMARY_BULLETS}}", narrative.executive_summary.bullets),
        ("{{KPI_COMMENTARY_BULLETS}}", narrative.kpi_commentary.bullets),
        ("{{SUPPLIER_INSIGHTS_BULLETS}}", narrative.supplier_insights.bullets),
        ("{{BUDGET_INSIGHTS_BULLETS}}", narrative.budget_insights.bullets),
    ]
    for slide in prs.slides:
        for token, bullets in bullet_slots:
            shp = _find_shape_with_token(slide, token)
            if shp:
                shp.text = shp.text.replace(token, "").strip()
                _set_bullets(shp, bullets)

    # Top suppliers table
    for slide in prs.slides:
        shp = _find_shape_with_token(slide, "{{TOP_SUPPLIERS_TABLE}}")
        if shp:
            suppliers = [(s.supplier, s.amount_net) for s in data.top_suppliers]
            _add_table_top_suppliers(slide, shp, suppliers)

    # Budget chart
    for slide in prs.slides:
        shp = _find_shape_with_token(slide, "{{BUDGET_CHART}}")
        if shp:
            cats = [(c.category_name, float(c.actual_net), float(c.budget)) for c in data.categories[:8]]
            _add_budget_chart(slide, shp, cats)

    # ============================================================
    # RISIKEN & MASSNAHMEN (Enterprise Feature)
    # ============================================================
    for slide in prs.slides:
        # Find and replace risk placeholder
        _replace_placeholder_with_bullets(
            slide, 
            "Risikoanalyse wird durch KI generiert",
            narrative.risks if narrative.risks else ["Keine signifikanten Risiken identifiziert."]
        )
        
        # Find and replace actions placeholder
        _replace_placeholder_with_bullets(
            slide,
            "Maßnahmen werden durch KI generiert", 
            narrative.actions if narrative.actions else ["Fortführung des aktuellen Kurses empfohlen."]
        )

    bio = BytesIO()
    prs.save(bio)
    return bio.getvalue()
